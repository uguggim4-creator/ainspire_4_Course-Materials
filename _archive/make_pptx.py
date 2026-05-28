# -*- coding: utf-8 -*-
# media-tools-deck.html 내용을 .pptx로 재구성 (텍스트 + 영상 삽입)
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DIR = r"f:\ainspire 4기 강의"
ASSETS = os.path.join(DIR, "demo-assets")

INK    = RGBColor(0x1A, 0x15, 0x30)
PURPLE = RGBColor(0x60, 0x20, 0xFF)
PURPLE_L = RGBColor(0xBB, 0xA4, 0xFF)
YELLOW = RGBColor(0xFF, 0xE6, 0x00)
DARK   = RGBColor(0x0D, 0x04, 0x20)
GRAY   = RGBColor(0x6B, 0x5D, 0x8A)
LINE   = RGBColor(0xE5, 0xDF, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Paperlogy"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def txt(s, left, top, width, height, text, size, color, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.1):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        run = p.add_run(); run.text = ln
        f = run.font; f.size = Pt(size); f.bold = bold
        f.color.rgb = color; f.name = font
    return tb

def step_label(s, text):
    txt(s, Inches(0.6), Inches(0.45), Inches(6), Inches(0.4), text, 13, PURPLE, bold=True)

def title(s, text, accent_part=None):
    tb = txt(s, Inches(0.6), Inches(0.85), Inches(8), Inches(1.1), text, 40, INK, bold=True)
    return tb

def divider(s, top=Inches(1.95)):
    ln = s.shapes.add_shape(1, Inches(0.62), top, Inches(0.55), Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb = PURPLE; ln.line.fill.background()
    ln.shadow.inherit = False

def info_block(s, left, top, width, label, body):
    bar = s.shapes.add_shape(1, left, top, Pt(3), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, left + Inches(0.18), top, width, Inches(0.3), label, 12, PURPLE, bold=True)
    txt(s, left + Inches(0.18), top + Inches(0.32), width, Inches(0.9), body, 13, INK, line_spacing=1.3)

def video_box(s, left, top, w, h, mp4, poster, label):
    txt(s, left, top - Inches(0.32), w, Inches(0.3), label, 12, PURPLE, bold=True)
    p = os.path.join(ASSETS, poster); v = os.path.join(ASSETS, mp4)
    if os.path.exists(v):
        s.shapes.add_movie(v, left, top, w, h,
                           poster_frame_image=p if os.path.exists(p) else None,
                           mime_type="video/mp4")

# ── S1 타이틀 ──
s = add_slide(DARK)
txt(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
    "AInspire 4기 · 미디어 도구 강의", 15, YELLOW, bold=True)
txt(s, Inches(0.9), Inches(2.5), Inches(11), Inches(1.8),
    "쇼츠 한 편을 만드는 여정", 54, WHITE, bold=True)
txt(s, Inches(0.9), Inches(4.4), Inches(11), Inches(0.5),
    "AI가 활용하는 미디어 도구 지도", 22, PURPLE_L)
txt(s, Inches(0.9), Inches(5.2), Inches(11), Inches(0.5),
    "yt-dlp    FFmpeg    Whisper    Remotion", 16, RGBColor(0xCC,0xCC,0xDD))

# ── S2 파이프라인 (4단계) ──
s = add_slide(WHITE)
step_label(s, "파이프라인이란?")
title(s, "단계가 이어지면 영상이 됩니다")
divider(s)
steps = [("Step 1","소스 구하기","yt-dlp"),("Step 2","자르고 다듬기","FFmpeg"),
         ("Step 3","자동 자막","Whisper"),("Step 4","조립·마무리","FFmpeg"),("출력","최종 파일",".mp4")]
bw = Inches(2.15); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(2.7)
for i,(st,nm,tool) in enumerate(steps):
    x = x0 + i*(bw+gap)
    box = s.shapes.add_shape(1, x, y0, bw, Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7,0xF3,0xFF) if i<4 else WHITE
    box.line.color.rgb = PURPLE if i<4 else LINE; box.line.width=Pt(1); box.shadow.inherit=False
    tf=box.text_frame; tf.word_wrap=True
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for j,(t,sz,c,b) in enumerate([(st,12,PURPLE,True),(nm,15,INK,True),(tool,14,PURPLE if i<4 else GRAY,True)]):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
txt(s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.0),
    "AI는 각 단계마다 다른 도구를 호출합니다. 명령어를 외울 필요 없이, 어떤 단계에서 어떤 도구가 왜 쓰이는지 이해하는 것이 핵심입니다.",
    14, INK, line_spacing=1.4)

# ── 도구 슬라이드 공통 ──
def tool_slide(label, ttl, intro, usage, req, res, v1, p1, v2, p2, l2):
    s = add_slide(WHITE)
    step_label(s, label)
    title(s, ttl)
    divider(s)
    LW = Inches(6.2)
    info_block(s, Inches(0.6), Inches(2.25), LW, "도구 소개", intro)
    info_block(s, Inches(0.6), Inches(3.35), LW, "활용 방식", usage)
    # 요청 박스
    txt(s, Inches(0.6), Inches(4.5), LW, Inches(0.3), "이렇게 요청하세요", 11, PURPLE, bold=True)
    rb = s.shapes.add_shape(1, Inches(0.6), Inches(4.82), LW, Inches(0.75))
    rb.fill.solid(); rb.fill.fore_color.rgb = PURPLE; rb.line.fill.background(); rb.shadow.inherit=False
    rtf=rb.text_frame; rtf.word_wrap=True; rtf.vertical_anchor=MSO_ANCHOR.MIDDLE
    rp=rtf.paragraphs[0]; rr=rp.add_run(); rr.text='"'+req+'"'
    rr.font.size=Pt(13); rr.font.color.rgb=WHITE; rr.font.name=FONT
    rtf.margin_left=Inches(0.2)
    # 결과
    txt(s, Inches(0.6), Inches(5.75), LW, Inches(0.3), "이런 결과를 받습니다", 11, PURPLE, bold=True)
    txt(s, Inches(0.6), Inches(6.05), LW, Inches(1.0), res, 12, INK, line_spacing=1.3)
    # 우측 영상 2개
    RX = Inches(7.1); RW = Inches(5.6); VH = Inches(2.35)
    video_box(s, RX, Inches(2.4), RW, VH, v1, p1, "1. 요청 → 실행")
    video_box(s, RX, Inches(5.15), RW, VH, v2, p2, l2)
    return s

tool_slide("Step 1 / yt-dlp", "yt-dlp — 소스 구하기",
  "yt-dlp는 YouTube를 비롯한 수백 개 온라인 플랫폼에서 영상·오디오를 내려받는 오픈소스 명령줄 도구입니다.",
  "영상 제작 파이프라인의 첫 단계입니다. 레퍼런스 클립, BGM, 원본 영상을 확보할 때 사용합니다.",
  "yt-dlp로 이 유튜브 영상 원본이랑 음성만 따로 받아줘",
  "원본 영상(mp4)과 오디오(mp3)가 저장되어 다음 단계의 작업 소스로 바로 활용됩니다.",
  "rec-ytdlp-flow.mp4","poster-rec-ytdlp-flow.png",
  "sample.mp4","poster-sample.png","2. 결과물 — 내려받은 원본 영상")

tool_slide("Step 2 / FFmpeg", "FFmpeg — 자르고 다듬기",
  "FFmpeg는 영상·오디오를 자르고 붙이고 변환하는 오픈소스 멀티미디어 처리 도구입니다.",
  "쇼츠·릴스는 9:16 세로가 기본입니다. 16:9 원본을 크롭·리사이즈하는 전처리 단계로 가장 빈번히 쓰입니다.",
  "FFmpeg로 이 영상 9:16 세로 비율로 잘라줘",
  "쇼츠 규격(1080x1920)으로 잘리고 크롭된 세로 영상 클립이 생성됩니다.",
  "rec-ffmpeg1-flow.mp4","poster-rec-ffmpeg1-flow.png",
  "vertical_9x16.mp4","poster-vertical_9x16.png","2. 결과물 — 9:16 세로 변환 클립")

tool_slide("Step 3 / Whisper", "Whisper — 자동 자막 생성",
  "Whisper는 OpenAI가 공개한 오픈소스 음성 인식(STT) 모델입니다. 99개 언어를 지원합니다.",
  "쇼츠에서 자막은 필수입니다. 음성을 자동으로 받아써 타임스탬프가 정확한 자막 파일을 만듭니다.",
  "Whisper로 이 영상 음성을 받아써서 자막 파일로 만들어줘",
  "타임코드가 들어간 자막(.srt) 파일이 자동 생성되어 바로 영상에 얹을 수 있습니다.",
  "rec-whisper-flow.mp4","poster-rec-whisper-flow.png",
  "rec-whisper-result.mp4","poster-rec-whisper-result.png","2. 결과물 — 생성된 자막 파일(.srt)")

tool_slide("Step 4 / FFmpeg", "FFmpeg — 조립·마무리",
  "전처리에 쓰인 FFmpeg가 조립 단계에서 다시 등장합니다. 자막 굽기, BGM 믹스, 최종 인코딩을 처리합니다.",
  "자막·BGM·편집 클립을 하나의 완성 파일로 합치는 파이프라인의 종착점입니다.",
  "편집한 영상에 자막이랑 BGM 합쳐서 최종 mp4로 뽑아줘",
  "자막과 음악이 합쳐진 업로드용 최종 영상 파일이 완성됩니다.",
  "rec-ffmpeg2-flow.mp4","poster-rec-ffmpeg2-flow.png",
  "final_shorts.mp4","poster-final_shorts.png","2. 결과물 — 자막 입힌 최종 쇼츠")

# ── 전체 흐름 ──
s = add_slide(WHITE)
step_label(s, "전체 흐름 다시보기")
title(s, "한눈에 보는 쇼츠 제작 파이프라인")
divider(s)
rows = [("Step 1","yt-dlp","YouTube 등에서 영상·오디오 파일을 내려받습니다."),
        ("Step 2","FFmpeg","구간 자르기, 포맷 변환, 9:16 세로 크롭을 수행합니다."),
        ("Step 3","Whisper","음성을 텍스트로 변환해 타임스탬프가 담긴 .srt 파일을 만듭니다."),
        ("Step 4","FFmpeg","자막 굽기·BGM 믹스 후 최종 mp4로 인코딩합니다.")]
y=Inches(2.5)
for st,tool,desc in rows:
    txt(s, Inches(0.8), y, Inches(1.3), Inches(0.4), st, 16, PURPLE, bold=True)
    txt(s, Inches(2.2), y, Inches(2.0), Inches(0.4), tool, 16, INK, bold=True)
    txt(s, Inches(4.3), y, Inches(8.2), Inches(0.4), desc, 14, INK)
    y += Inches(0.85)
txt(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
    "각 단계는 독립적으로 실행되며, AI는 상황에 따라 필요한 단계만 선택적으로 호출합니다.", 13, GRAY)

# ── 용어집 ──
s = add_slide(WHITE)
step_label(s, "부록")
title(s, "알아두면 좋은 용어들")
divider(s)
terms = [("코덱 vs 컨테이너","코덱은 영상을 압축하는 방식(H.264 등), 컨테이너는 그것을 담는 파일 형식(mp4 등)입니다."),
         ("비트레이트","1초에 담기는 데이터 양. 높을수록 화질이 좋지만 용량이 커집니다."),
         ("fps","초당 프레임 수. 30fps는 1초에 30장의 그림이 지나갑니다."),
         ("종횡비","화면의 가로:세로 비율. 가로 16:9, 쇼츠는 세로 9:16입니다."),
         ("STT","음성을 텍스트로 바꾸는 기술(Speech-to-Text). Whisper가 대표적입니다."),
         ("렌더링","코드·데이터를 실제 영상 프레임으로 그려내는 과정입니다.")]
y=Inches(2.4); col_x=[Inches(0.8), Inches(6.9)]
for i,(t,d) in enumerate(terms):
    cx = col_x[i%2]; cy = y + (i//2)*Inches(1.5)
    txt(s, cx, cy, Inches(5.5), Inches(0.4), t, 17, PURPLE, bold=True)
    txt(s, cx, cy+Inches(0.42), Inches(5.5), Inches(0.9), d, 13, INK, line_spacing=1.3)

# ── Remotion (보너스) ──
s = add_slide(WHITE)
step_label(s, "보너스 / Remotion")
title(s, "Remotion — 코드로 영상 만들기")
divider(s)
LW=Inches(6.2)
info_block(s, Inches(0.6), Inches(2.25), LW, "도구 소개",
    "Remotion은 React 컴포넌트로 영상의 각 프레임을 정의하고 mp4로 렌더링하는 프로그래매틱 영상 생성 프레임워크입니다.")
info_block(s, Inches(0.6), Inches(3.5), LW, "활용 방식",
    "자막 타이밍·조회수 같은 데이터를 코드로 받아 애니메이션 영상으로 자동 변환합니다. 같은 템플릿으로 수십 편을 일괄 제작할 때 특히 유용합니다.")
txt(s, Inches(0.6), Inches(4.9), LW, Inches(0.3), "이렇게 요청하세요", 11, PURPLE, bold=True)
rb=s.shapes.add_shape(1, Inches(0.6), Inches(5.22), LW, Inches(0.75))
rb.fill.solid(); rb.fill.fore_color.rgb=PURPLE; rb.line.fill.background(); rb.shadow.inherit=False
rtf=rb.text_frame; rtf.word_wrap=True; rtf.vertical_anchor=MSO_ANCHOR.MIDDLE; rtf.margin_left=Inches(0.2)
rr=rtf.paragraphs[0].add_run(); rr.text='"Remotion으로 조회수 데이터를 막대그래프 영상으로 자동 생성해줘"'
rr.font.size=Pt(12); rr.font.color.rgb=WHITE; rr.font.name=FONT
video_box(s, Inches(7.1), Inches(2.6), Inches(5.6), Inches(4.0),
          "rec-remotion-result.mp4","poster-rec-remotion-result.png","렌더링 결과 — 실제 Remotion 영상")

# ── 마무리 ──
s = add_slide(DARK)
txt(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.5),
    "도구는 손, 머리는 당신", 46, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.2), Inches(11.3), Inches(1.0),
    "AI는 이 도구들을 순서대로 불러줍니다.\n당신은 무엇을 만들지만 정하면 됩니다.", 22, PURPLE_L,
    align=PP_ALIGN.CENTER, line_spacing=1.4)

out = os.path.join(DIR, "media-tools-deck.pptx")
prs.save(out)
print("저장 완료:", out)
print("슬라이드 수:", len(prs.slides._sldIdLst))
